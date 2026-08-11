from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

from app.export.contracts import ExportPayload


def export_pptx(payload: ExportPayload, output_path: str | Path) -> Path:
    path = Path(output_path)
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(12), Inches(0.7))
    title.text_frame.text = payload.title
    title.text_frame.paragraphs[0].font.size = Pt(28)
    title.text_frame.paragraphs[0].font.bold = True

    subtitle = slide.shapes.add_textbox(Inches(0.7), Inches(1.15), Inches(12), Inches(0.45))
    subtitle.text_frame.text = f"{payload.subtitle} · {payload.period}"
    subtitle.text_frame.paragraphs[0].font.size = Pt(16)

    labels = list(payload.metrics.items())[:4]
    positions = [(0.8, 2.0), (3.9, 2.0), (7.0, 2.0), (10.1, 2.0)]
    for (label, value), (x, y) in zip(labels, positions):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(2.4), Inches(1.35))
        box.text_frame.text = f"{label}\n{value}"
        box.text_frame.paragraphs[0].font.size = Pt(13)
        box.text_frame.paragraphs[1].font.size = Pt(25)
        box.text_frame.paragraphs[1].font.bold = True

    status = slide.shapes.add_textbox(Inches(0.8), Inches(4.0), Inches(5), Inches(0.6))
    status.text_frame.text = f"Estado: {payload.status}"
    status.text_frame.paragraphs[0].font.size = Pt(20)
    status.text_frame.paragraphs[0].font.bold = True

    alert_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.8), Inches(11.8), Inches(1.2))
    alert_text = "\n".join(payload.alerts) if payload.alerts else "Sin alertas"
    alert_box.text_frame.text = f"Alertas\n{alert_text}"
    alert_box.text_frame.paragraphs[0].font.size = Pt(16)
    alert_box.text_frame.paragraphs[0].font.bold = True
    if len(alert_box.text_frame.paragraphs) > 1:
        alert_box.text_frame.paragraphs[1].font.size = Pt(14)

    evidence = slide.shapes.add_textbox(Inches(0.8), Inches(6.55), Inches(11.8), Inches(0.35))
    evidence.text_frame.text = f"Evidencia: {payload.evidence_id}"
    evidence.text_frame.paragraphs[0].font.size = Pt(10)

    presentation.save(path)
    return path
