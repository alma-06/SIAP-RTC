from __future__ import annotations

from dataclasses import asdict
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

from app.export.executive_package import build_executive_summary
from app.pipeline.integrated import IntegratedPipelineResult


def _add_bullets(slide, items: list[str]) -> None:
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.2))
    frame = box.text_frame
    frame.clear()
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.font.size = Pt(20)
        paragraph.space_after = Pt(10)


def write_pptx_package(result: IntegratedPipelineResult, path: str | Path) -> None:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    summary = build_executive_summary(result)

    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    slide.shapes.title.text = "SIAP-RTC — Resumen ejecutivo"
    slide.placeholders[1].text = f"Periodo: {', '.join(summary['periods'])} | Evidencia: {summary['evidence_id']}"

    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Resultados del procesamiento"
    _add_bullets(slide, [
        f"Registros ingeridos: {summary['records_ingested']}",
        f"Registros normalizados: {summary['records_normalized']}",
        f"Registros conservados: {summary['records_kept']}",
        f"Duplicados eliminados: {summary['duplicates_removed']}",
    ])

    if result.indicators is not None:
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = "Indicadores de conciliación"
        indicators = asdict(result.indicators)
        _add_bullets(slide, [
            f"Periodo anterior: {indicators['previous_count']}",
            f"Periodo actual: {indicators['current_count']}",
            f"Permanencias: {indicators['unchanged_count']}",
            f"Altas: {indicators['added_count']} | Bajas: {indicators['removed_count']} | Modificados: {indicators['modified_count']}",
            f"Variación neta: {indicators['net_change']}",
        ])

    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Consideraciones metodológicas"
    _add_bullets(slide, [*summary['warnings'], "La pauta RTC acredita programación publicada; no constituye por sí misma prueba de transmisión efectiva."])

    Path(path).parent.mkdir(parents=True, exist_ok=True)
    presentation.save(path)
