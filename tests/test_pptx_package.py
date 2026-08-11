from openpyxl import Workbook
from pptx import Presentation

from app.export.pptx_package import write_pptx_package
from app.pipeline.integrated import run_integrated_period

HEADERS = ["Pauta de transmisión", "Estado", "Tiempo Fiscal", "Canal Base", "Orden", "Fecha", "Dependencia CAM. SEN.", "Clave", "Campaña", "Versión"]


def test_write_pptx_package_creates_executive_slides(tmp_path) -> None:
    source = tmp_path / "rtc.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.append(HEADERS)
    sheet.append(["Pauta", "Programada", "00:00:30", "AM", "O1", "11/08/2026", "CAM. SEN.", "C1", "Campaña", "V1"])
    workbook.save(source)

    result = run_integrated_period([source], "2026-Q2")
    output = tmp_path / "Presentacion_Ejecutiva.pptx"
    write_pptx_package(result, output)
    presentation = Presentation(output)
    assert len(presentation.slides) == 3
    assert presentation.slides[0].shapes.title.text == "SIAP-RTC — Resumen ejecutivo"
