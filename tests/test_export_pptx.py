from pptx import Presentation

from app.dashboard.fact_sheet import build_fact_sheet
from app.dashboard.model import DashboardPeriod
from app.export.contracts import build_export_payload
from app.export.pptx import export_pptx


def test_export_pptx_creates_executive_slide(tmp_path) -> None:
    period = DashboardPeriod("2026-Q2", 100, 10, 5, 80, 5, 0.80, 0.05, 0.10, 0.05, True, (), "EV-01")
    payload = build_export_payload(build_fact_sheet(period))
    output = export_pptx(payload, tmp_path / "siap_rtc.pptx")
    presentation = Presentation(output)
    assert len(presentation.slides) == 1
    text = "\n".join(shape.text for shape in presentation.slides[0].shapes if hasattr(shape, "text"))
    assert "SIAP-RTC" in text
    assert "2026-Q2" in text
    assert "Total comparado" in text
    assert "EV-01" in text
